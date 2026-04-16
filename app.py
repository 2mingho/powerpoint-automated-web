import os
import pandas as pd
import uuid
import zipfile
import shutil
import json
import threading
import random
from datetime import datetime, timedelta
import functools
import click
from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for, abort, after_this_request, flash, session
from flask_login import current_user, login_required
from services.classifier import classify_mentions
from services.file_loader import detect_format, read_full_as_tsv
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from babel.dates import format_datetime
from sqlalchemy import inspect, text, event
from dotenv import load_dotenv
from flask_talisman import Talisman

from blueprints.auth import auth
from blueprints.admin import admin_bp, log_activity
from blueprints.tasks import tasks_bp


from extensions import db, login_manager, csrf, limiter
from models import User, Report, ActivityLog, ClassificationPreset, Task, TempArtifact
from services import calculation as report
from services.groq_analysis import construir_prompt, llamar_groq, extraer_json, formatear_analisis_social_listening
from pptx_builder import engine as ppt_engine
from pptx_builder import native_charts
from pptx_builder.engine import set_text_style
from services.csv_analysis import analyze_csv, generate_summary_csv

# Load environment variables
load_dotenv()


def _env_bool(name, default=False):
    raw = os.environ.get(name)
    if raw is None:
        return default
    return raw.strip().lower() in {'1', 'true', 'yes', 'on'}


def _env_int(name, default):
    raw = os.environ.get(name)
    if raw is None:
        return default
    try:
        return int(raw)
    except (TypeError, ValueError):
        return default


def _env_float(name, default):
    raw = os.environ.get(name)
    if raw is None:
        return default
    try:
        return float(raw)
    except (TypeError, ValueError):
        return default


def _resolve_database_uri():
    uri = os.environ.get('DATABASE_URL') or os.environ.get('SQLALCHEMY_DATABASE_URI')
    if uri:
        if uri.startswith('postgres://'):
            return 'postgresql://' + uri[len('postgres://'):]
        return uri
    return 'sqlite:///users.db'


def _is_production_mode():
    return (
        os.environ.get('FLASK_ENV') == 'production'
        or _env_bool('RENDER', False)
        or _env_bool('FORCE_PRODUCTION_MODE', False)
    )


app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'scratch'

# Security configuration
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY')
if not app.config['SECRET_KEY']:
    raise RuntimeError("SECRET_KEY must be set in .env file. Generate one with: python -c 'import secrets; print(secrets.token_hex(32))'")

app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # 200 MB upload limit
app.config['SQLALCHEMY_DATABASE_URI'] = _resolve_database_uri()
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['ALLOW_SELF_REGISTRATION'] = _env_bool('ALLOW_SELF_REGISTRATION', False)

if app.config['SQLALCHEMY_DATABASE_URI'].startswith('postgresql://'):
    app.config['SQLALCHEMY_ENGINE_OPTIONS'] = {
        'pool_pre_ping': True,
        'pool_recycle': 300,
    }

# Session security
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(hours=8)
# Only set Secure cookie when not in debug/local mode
if _is_production_mode():
    app.config['SESSION_COOKIE_SECURE'] = True

# ─────────────────────────────────────────────────────────────
# Initialize extensions
# ─────────────────────────────────────────────────────────────
db.init_app(app)
login_manager.init_app(app)
csrf.init_app(app)
limiter.init_app(app)

# Security headers via Talisman (CSP, HSTS, X-Frame-Options)
# Using 'unsafe-inline' for scripts/styles since templates use inline code extensively
# NOTE: Do NOT use content_security_policy_nonce_in — it causes browsers to ignore 'unsafe-inline'
Talisman(app,
         force_https=_is_production_mode(),
         content_security_policy={
             'default-src': "'self'",
             'script-src': ["'self'", "'unsafe-inline'", "'unsafe-eval'", "https://cdn.jsdelivr.net", "https://cdnjs.cloudflare.com"],
             'style-src':  ["'self'", "'unsafe-inline'", "https://fonts.googleapis.com", "https://cdnjs.cloudflare.com"],
             'font-src':   ["'self'", "https://fonts.gstatic.com", "https://cdnjs.cloudflare.com"],
             'img-src':    ["'self'", "data:"],
             'connect-src': "'self'",
         })

# SQLite WAL mode for better concurrent access
from sqlalchemy import Engine as _Engine

@event.listens_for(_Engine, "connect")
def _set_sqlite_pragma(dbapi_conn, connection_record):
    import sqlite3
    if isinstance(dbapi_conn, sqlite3.Connection):
        cursor = dbapi_conn.cursor()
        cursor.execute("PRAGMA journal_mode=WAL")
        cursor.execute("PRAGMA busy_timeout=5000")
        cursor.close()

# Registrar blueprints
app.register_blueprint(auth)
app.register_blueprint(admin_bp)
app.register_blueprint(tasks_bp)

# ─────────────────────────────────────────────────────────────
# Force-logout check (session kick feature)
# ─────────────────────────────────────────────────────────────
from flask_login import logout_user

@app.before_request
def check_force_logout():
    """If admin has flagged this user for forced logout, log them out immediately."""
    session_user_id = session.get('_user_id')
    if session_user_id:
        try:
            session_user = db.session.get(User, int(session_user_id))
        except (TypeError, ValueError):
            session_user = None
        if session_user and not session_user.is_active:
            logout_user()
            session.pop('_user_id', None)
            session.pop('_fresh', None)
            flash('Tu cuenta está inactiva. Contacta al administrador.', 'warning')
            return redirect(url_for('auth.login'))

    if current_user.is_authenticated and not current_user.is_active:
        logout_user()
        flash('Tu cuenta está inactiva. Contacta al administrador.', 'warning')
        return redirect(url_for('auth.login'))

    if current_user.is_authenticated and getattr(current_user, 'force_logout', False):
        current_user.force_logout = False
        db.session.commit()
        logout_user()
        flash('Tu sesion ha sido terminada por un administrador.', 'warning')
        return redirect(url_for('auth.login'))

# ─────────────────────────────────────────────────────────────
# Automatic Request Logging (after_request)
# ─────────────────────────────────────────────────────────────

# Endpoints that already log manually — skip auto-log to avoid duplicates
_MANUALLY_LOGGED = {
    'index', 'clasificacion', 'download_file', 'download_classified',
    'clasificacion_finalize', 'analisis_csv', 'auth.login', 'auth.logout',
    'auth.register', 'union_archivos', 'union_detect', 'union_merge',
    'union_download',
}

DEFAULT_ENABLE_PAGE_VIEW_LOGS = not _is_production_mode()
ENABLE_PAGE_VIEW_LOGS = _env_bool('ENABLE_PAGE_VIEW_LOGS', DEFAULT_ENABLE_PAGE_VIEW_LOGS)
PAGE_VIEW_LOG_SAMPLE_RATE = max(0.0, min(1.0, _env_float(
    'PAGE_VIEW_LOG_SAMPLE_RATE',
    1.0 if ENABLE_PAGE_VIEW_LOGS else 0.0,
)))

@app.after_request
def auto_log_request(response):
    """Automatically log successful authenticated GET page-views."""
    if not ENABLE_PAGE_VIEW_LOGS:
        return response
    if PAGE_VIEW_LOG_SAMPLE_RATE <= 0:
        return response
    if PAGE_VIEW_LOG_SAMPLE_RATE < 1.0 and random.random() > PAGE_VIEW_LOG_SAMPLE_RATE:
        return response

    if (current_user.is_authenticated
            and response.status_code < 400
            and request.endpoint
            and not request.endpoint.startswith('static')
            and not request.endpoint.startswith('admin.')
            and request.endpoint not in _MANUALLY_LOGGED
            and not request.is_json
            and request.method == 'GET'):
        try:
            log_activity('page_view', f'{request.method} {request.endpoint}')
        except Exception:
            pass
    return response


if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])


# ✅ NUEVO: aseguramos esquema en la misma DB que usa la app
def ensure_default_admin():
    """Create default admin if none exists (idempotent)."""
    admin_email = os.environ.get('ADMIN_EMAIL', 'admin@dataintel.com')
    admin_password = os.environ.get('ADMIN_PASSWORD', 'Admin2024!')
    admin_username = os.environ.get('ADMIN_USERNAME', 'admin')

    if _is_production_mode() and admin_password == 'Admin2024!':
        raise RuntimeError(
            "Startup blocked: ADMIN_PASSWORD is using insecure default value in production. "
            "Set a strong ADMIN_PASSWORD environment variable."
        )

    existing_admin = User.query.filter_by(role='admin').first()
    if existing_admin:
        return existing_admin

    existing_by_email = User.query.filter_by(email=admin_email).first()
    if existing_by_email:
        if existing_by_email.role != 'admin':
            existing_by_email.role = 'admin'
            existing_by_email.is_active = True
            db.session.commit()
        return existing_by_email

    admin_user = User(
        username=admin_username,
        email=admin_email,
        password=generate_password_hash(admin_password, method='scrypt'),
        role='admin',
        is_active=True,
    )
    db.session.add(admin_user)
    db.session.commit()
    app.logger.info(f"[startup] Admin created: {admin_username} ({admin_email})")
    return admin_user


def ensure_reports_schema():
    """
    Garantiza que la tabla reports exista y tenga la columna template_name.
    Se ejecuta en el arranque de la app, usando la MISMA conexión de SQLAlchemy.
    """
    with app.app_context():
        # Crea tablas faltantes (idempotente)
        db.create_all()

        insp = inspect(db.engine)
        tables = set(insp.get_table_names())
        if 'reports' in tables:
            cols = {c['name'] for c in insp.get_columns('reports')}
            if 'template_name' not in cols:
                try:
                    db.session.execute(text("ALTER TABLE reports ADD COLUMN template_name TEXT"))
                    db.session.commit()
                except Exception as e:
                    print(f"[ensure_reports_schema] Aviso al agregar columna template_name: {e}")

        # Migrate users table for session control + area columns
        if 'users' in tables:
            user_cols = {c['name'] for c in insp.get_columns('users')}
            new_user_cols = {
                'session_token': 'VARCHAR(64)',
                'force_logout': 'BOOLEAN DEFAULT 0',
                'area_id': 'INTEGER REFERENCES areas(id)',
            }
            for col_name, col_type in new_user_cols.items():
                if col_name not in user_cols:
                    try:
                        db.session.execute(text(f"ALTER TABLE users ADD COLUMN {col_name} {col_type}"))
                        db.session.commit()
                        print(f"[migration] Added column users.{col_name}")
                    except Exception as e:
                        print(f"[migration] Aviso al agregar users.{col_name}: {e}")

        if 'tasks' in tables:
            task_cols = {c['name'] for c in insp.get_columns('tasks')}
            new_task_cols = {
                'start_date': 'DATE',
                'end_date': 'DATE',
                'directorate': 'VARCHAR(255)',
                'requested_by': 'VARCHAR(255)',
                'budget_type': 'VARCHAR(255)',
            }
            for col_name, col_type in new_task_cols.items():
                if col_name not in task_cols:
                    try:
                        db.session.execute(text(f"ALTER TABLE tasks ADD COLUMN {col_name} {col_type}"))
                        db.session.commit()
                        print(f"[migration] Added column tasks.{col_name}")
                    except Exception as e:
                        print(f"[migration] Aviso al agregar tasks.{col_name}: {e}")

        ensure_default_admin()


ACTIVITY_LOG_RETENTION_DAYS = max(1, _env_int('ACTIVITY_LOG_RETENTION_DAYS', 90))
ACTIVITY_LOG_MAX_ROWS = max(1000, _env_int('ACTIVITY_LOG_MAX_ROWS', 100000))
REPORT_METADATA_RETENTION_DAYS = max(1, _env_int('REPORT_METADATA_RETENTION_DAYS', 180))


def prune_activity_logs(retention_days=ACTIVITY_LOG_RETENTION_DAYS, max_rows=ACTIVITY_LOG_MAX_ROWS):
    """Prune old activity logs by retention and hard row cap."""
    deleted_by_age = 0
    deleted_by_cap = 0

    cutoff = datetime.utcnow() - timedelta(days=retention_days)
    deleted_by_age = (
        ActivityLog.query
        .filter(ActivityLog.timestamp < cutoff)
        .delete(synchronize_session=False)
    )
    db.session.commit()

    total_rows = ActivityLog.query.count()
    overflow = max(0, total_rows - max_rows)
    if overflow > 0:
        ids_to_delete = [
            row.id
            for row in (
                ActivityLog.query
                .with_entities(ActivityLog.id)
                .order_by(ActivityLog.timestamp.asc())
                .limit(overflow)
                .all()
            )
        ]
        if ids_to_delete:
            deleted_by_cap = (
                ActivityLog.query
                .filter(ActivityLog.id.in_(ids_to_delete))
                .delete(synchronize_session=False)
            )
            db.session.commit()

    return {
        'deleted_by_age': deleted_by_age,
        'deleted_by_cap': deleted_by_cap,
        'remaining_rows': ActivityLog.query.count(),
    }


def prune_report_metadata(retention_days=REPORT_METADATA_RETENTION_DAYS):
    """Prune report metadata rows older than retention window."""
    cutoff = datetime.utcnow() - timedelta(days=retention_days)
    deleted = (
        Report.query
        .filter(Report.created_at < cutoff)
        .delete(synchronize_session=False)
    )
    db.session.commit()
    return {'deleted_rows': deleted}


def _scratch_root_abs():
    return os.path.abspath(app.config['UPLOAD_FOLDER'])


def _scratch_path(filename):
    return os.path.join(app.config['UPLOAD_FOLDER'], filename)


def _register_temp_artifact(kind, file_id, storage_name, user_id=None):
    """Create/update ownership metadata for temporary downloadable files."""
    safe_file_id = secure_filename(file_id)
    if not safe_file_id:
        raise ValueError('file_id invalido')

    owner_id = user_id or current_user.id
    artifact = TempArtifact.query.filter_by(kind=kind, file_id=safe_file_id).first()
    if artifact is None:
        artifact = TempArtifact(
            kind=kind,
            file_id=safe_file_id,
            storage_name=storage_name,
            user_id=owner_id,
        )
        db.session.add(artifact)
    else:
        artifact.storage_name = storage_name
        artifact.user_id = owner_id

    db.session.commit()
    return artifact


def _get_owned_artifact_or_403(kind, file_id):
    """Load artifact metadata and enforce owner-or-admin access."""
    safe_file_id = secure_filename(file_id)
    artifact = TempArtifact.query.filter_by(kind=kind, file_id=safe_file_id).first()
    if artifact is None:
        abort(404)
    if not current_user.is_admin and artifact.user_id != current_user.id:
        abort(403)
    return artifact


def prune_database_storage():
    """Run all DB pruning tasks and return stats."""
    logs_stats = prune_activity_logs()
    reports_stats = prune_report_metadata()
    return {
        'activity_logs': logs_stats,
        'reports': reports_stats,
    }


@app.cli.command('maintenance-prune')
def maintenance_prune_command():
    """Prune DB metadata tables to control storage usage."""
    with app.app_context():
        stats = prune_database_storage()
    click.echo(f"Activity logs pruned: {stats['activity_logs']}")
    click.echo(f"Reports metadata pruned: {stats['reports']}")


# Ejecutar guardado de esquema al iniciar
ensure_reports_schema()

if _env_bool('RUN_STARTUP_MAINTENANCE', True):
    try:
        with app.app_context():
            stats = prune_database_storage()
        app.logger.info(f"[startup] DB pruning completed: {stats}")
    except Exception as e:
        app.logger.warning(f"[startup] DB pruning warning: {e}")

try:
    with app.app_context():
        app.logger.info(f"[startup] Database URL: {db.engine.url.render_as_string(hide_password=True)}")
        app.logger.info(f"[startup] Admin users: {User.query.filter_by(role='admin').count()}")
except Exception as e:
    app.logger.warning(f"[startup] DB diagnostics warning: {e}")


# ─────────────────────────────────────────────────────────────
# Utilidades para plantillas
DEFAULT_TEMPLATE_FILENAME = "Reporte_plantilla.pptx"
TEMPLATES_DIR = "powerpoints"


def get_available_templates():
    """Lista segura de nombres de archivo .pptx dentro de powerpoints/"""
    try:
        files = [
            f for f in os.listdir(TEMPLATES_DIR)
            if os.path.isfile(os.path.join(TEMPLATES_DIR, f)) and f.lower().endswith(".pptx")
        ]
        return sorted(files)
    except FileNotFoundError:
        return []


def template_path_from_name(template_name):
    """Construye la ruta absoluta segura a la plantilla."""
    safe_name = os.path.basename(template_name)  # evita traversal
    return os.path.join(TEMPLATES_DIR, safe_name)


def clean_scratch_folder():
    """Clean only old files (>1 hour) to prevent race conditions.
    Runs in a background thread — never blocks user requests."""
    folder = app.config['UPLOAD_FOLDER']
    try:
        import time
        current_time = time.time()
        one_hour_ago = current_time - 3600
        
        for file in os.listdir(folder):
            file_path = os.path.join(folder, file)
            try:
                # Only delete files/folders older than 1 hour
                if os.path.getmtime(file_path) < one_hour_ago:
                    if os.path.isfile(file_path):
                        os.unlink(file_path)
                    elif os.path.isdir(file_path):
                        shutil.rmtree(file_path)
            except Exception:
                continue
    except Exception as e:
        app.logger.error(f"Error al limpiar la carpeta scratch: {e}")


def _schedule_background_cleanup():
    """Start a background thread that cleans scratch/ every 30 minutes."""
    def _run():
        while True:
            import time
            time.sleep(1800)  # 30 minutes
            with app.app_context():
                clean_scratch_folder()
    t = threading.Thread(target=_run, daemon=True)
    t.start()

_schedule_background_cleanup()


# ─────────────────────────────────────────────────────────────
# Tool access decorator
# ─────────────────────────────────────────────────────────────

def _is_ajax():
    """Return True if the request is an AJAX/fetch call expecting JSON."""
    return (
        request.is_json
        or request.headers.get('X-Requested-With') == 'XMLHttpRequest'
        or 'application/json' in request.headers.get('Accept', '')
        or request.headers.get('Content-Type', '').startswith('multipart')  # FormData fetch
    )


def tool_required(tool_key):
    """Decorator: blocks access if user lacks permission for *tool_key*.
    For AJAX/fetch requests returns JSON errors instead of HTML redirects.
    """
    def wrapper(f):
        @functools.wraps(f)
        def decorated(*args, **kwargs):
            if not current_user.is_authenticated:
                if _is_ajax():
                    from flask import jsonify as _jsonify
                    return _jsonify({"success": False, "error": "Sesion expirada. Por favor recarga la pagina e inicia sesion nuevamente."}), 401
                return redirect(url_for('auth.login'))
            if not current_user.has_tool_access(tool_key):
                if _is_ajax():
                    from flask import jsonify as _jsonify
                    return _jsonify({"success": False, "error": "No tienes permiso para acceder a esta herramienta."}), 403
                flash('No tienes permiso para acceder a esta herramienta.', 'error')
                return redirect(url_for('menu'))
            return f(*args, **kwargs)
        return decorated
    return wrapper


# Inject has_tool_access into templates
@app.context_processor
def inject_tool_access():
    def _has_tool_access(tool_key):
        if current_user.is_authenticated:
            return current_user.has_tool_access(tool_key)
        return False
    return dict(has_tool_access=_has_tool_access)


@app.route('/menu')
@login_required
def menu():
    return render_template('menu.html')


@app.route('/', methods=['GET', 'POST'])
@tool_required('reports')
def index():
    available_templates = get_available_templates()
    default_template = DEFAULT_TEMPLATE_FILENAME if DEFAULT_TEMPLATE_FILENAME in available_templates else (available_templates[0] if available_templates else None)

    if request.method == 'POST':
        csv_file = request.files.get('csv_file')
        wordcloud_file = request.files.get('wordcloud_file')
        report_title = request.form.get('report_title', '').strip()
        description = request.form.get('description', '').strip()
        selected_template = request.form.get('template_name', default_template)
        fallback_default = bool(request.form.get('fallback_default'))

        if not csv_file or not csv_file.filename.endswith('.csv'):
            return redirect(url_for('error_archivo_invalido'))

        # Validar plantilla seleccionada
        if not available_templates:
            flash("No hay plantillas disponibles en el servidor. Contacta al administrador.", "error")
            return render_template('index.html',
                                   available_templates=[],
                                   default_template=None)

        if selected_template not in available_templates:
            if fallback_default and default_template:
                flash(f"La plantilla seleccionada no existe. Se usará la predeterminada: {default_template}.", "warning")
                selected_template = default_template
            else:
                flash("La plantilla seleccionada no existe. Por favor elige otra o marca 'Usar la plantilla predeterminada'.", "error")
                return render_template('index.html',
                                       available_templates=available_templates,
                                       default_template=default_template)

        unique_id = uuid.uuid4().hex[:6]
        csv_filename = secure_filename(csv_file.filename)
        csv_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{csv_filename}")
        csv_file.save(csv_path)

        wordcloud_path = None
        if wordcloud_file and wordcloud_file.filename.endswith('.png'):
            wordcloud_path = os.path.join(app.config['UPLOAD_FOLDER'], 'Wordcloud.png')
            wordcloud_file.save(wordcloud_path)

        try:
            zip_path, missing_fields, used_template = process_report(
                csv_path=csv_path,
                wordcloud_path=wordcloud_path,
                unique_id=unique_id,
                template_filename=selected_template,
                report_title=report_title,
                description=description
            )
            log_activity('generate_report', f'Reporte generado: {report_title or csv_filename} (plantilla: {selected_template})')
        except Exception as e:
            app.logger.error(f"Error generando el reporte: {e}")
            abort(500)

        if missing_fields:
            faltantes = ", ".join(sorted(set(missing_fields)))
            flash(f"Advertencia: La plantilla '{used_template}' no contiene algunos campos esperados y fueron omitidos: {faltantes}.", "warning")

        zip_filename = os.path.basename(zip_path)
        file_size_mb = round(os.path.getsize(zip_path) / (1024 * 1024), 2)
        current_time = datetime.now()
        formatted_datetime = format_datetime(current_time, "d 'de' MMMM, yyyy - HH:mm", locale='es')

        return render_template(
            'download.html',
            zip_path=zip_filename,
            file_size=file_size_mb,
            formatted_datetime=formatted_datetime,
            template_used=used_template
        )

    # GET
    return render_template('index.html',
                           available_templates=available_templates,
                           default_template=default_template)


def process_report(csv_path, wordcloud_path, unique_id, template_filename, report_title=None, description=None):
    """
    Genera el zip del reporte usando la plantilla indicada.
    Devuelve: (zip_path, missing_fields_list, used_template_filename)
    """
    missing_fields = []

    # Carga y limpieza
    df_cleaned = report.load_and_clean_data(csv_path)
    df_cleaned['Influencer'] = df_cleaned.apply(report.update_influencer, axis=1)
    df_cleaned['Sentiment'] = df_cleaned.apply(report.update_sentiment, axis=1)

    total_mentions, count_of_authors, estimated_reach = report.calculate_summary_metrics(df_cleaned)
    processed_csv_path = report.save_cleaned_csv(df_cleaned, csv_path, unique_id)

    solo_fecha = request.form.get('solo_fecha') is not None
    evolution_data = report.get_evolution_data(df_cleaned, use_date_only=not solo_fecha)
    sentiment_data = report.get_sentiment_data(df_cleaned)

    platform_counts, _ = report.distribucion_plataforma(df_cleaned)
    top_sentences = report.get_top_hit_sentences(df_cleaned)
    top_influencers_prensa = report.get_top_influencers(df_cleaned, 'Prensa Digital', sort_by='Posts')
    top_influencers_redes_posts = report.get_top_influencers(df_cleaned, 'Redes Sociales', sort_by='Posts', include_source=True)
    top_influencers_redes_reach = report.get_top_influencers(df_cleaned, 'Redes Sociales', sort_by='Max Reach')

    current_date = datetime.now().strftime('%d-%b-%Y')
    client_name = report_title if report_title else os.path.basename(csv_path).split()[0]

    # Abrir plantilla seleccionada
    tpl_path = template_path_from_name(template_filename)
    if not os.path.isfile(tpl_path):
        raise FileNotFoundError(f"Plantilla no encontrada: {tpl_path}")

    prs = Presentation(tpl_path)

    # --- OPTIMIZED: Single-pass placeholder indexing (P3) ---
    # Build index once instead of scanning slides multiple times
    placeholder_index = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text.strip():
                    placeholder_index[shape.text.strip()] = (slide, shape)
            except Exception:
                continue

    def find_shape_for_key(key):
        """Fast lookup using pre-built index"""
        return placeholder_index.get(key, (None, None))

    # Reemplazo de textos genéricos
    text_mapping = {
        "REPORT_CLIENT": client_name,
        "REPORT_DATE": current_date,
        "NUMB_MENTIONS": str(total_mentions),
        "NUMB_ACTORS": str(count_of_authors),
        "EST_REACH": estimated_reach
    }

    # Apply text replacements using index
    found_text_keys = set()
    for key, value in text_mapping.items():
        slide, shape = find_shape_for_key(key)
        if shape:
            try:
                # Custom color for REPORT_DATE (white)
                text_color = RGBColor(255, 255, 255) if key == "REPORT_DATE" else RGBColor(0, 0, 0)
                set_text_style(shape, str(value), 'Effra Heavy', Pt(28), 
                               False if key not in ("NUMB_MENTIONS", "NUMB_ACTORS", "EST_REACH") else True,
                               color=text_color)
                found_text_keys.add(key)
            except Exception:
                pass
    
    for key in text_mapping.keys():
        if key not in found_text_keys:
            missing_fields.append(key)

    # Charts / imágenes: buscar placeholder y añadir imagen en la ubicación del placeholder
    def place_image_at_placeholder(key, image_path, default_size=None):
        slide, shape = find_shape_for_key(key)
        if slide and shape:
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                # eliminar texto para evitar superposición
                try:
                    shape.text = ""
                except Exception:
                    pass
                if default_size:
                    width, height = default_size
                slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                return True
            except Exception:
                return False
        return False

    # Conversación (native line chart)
    conv_slide, conv_shape = find_shape_for_key('CONVERSATION_CHART')
    if conv_slide and conv_shape:
        try:
            native_charts.add_native_line_chart(
                conv_slide, conv_shape,
                evolution_data['labels'], evolution_data['values'],
                width=Inches(9.07), height=Inches(5.15)
            )
        except Exception:
            missing_fields.append('CONVERSATION_CHART')
    else:
        missing_fields.append('CONVERSATION_CHART')

    # Sentiment pie (native pie chart)
    sent_slide, sent_shape = find_shape_for_key('SENTIMENT_PIE')
    if sent_slide and sent_shape:
        try:
            native_charts.add_native_pie_chart(sent_slide, sent_shape, sentiment_data, width=Inches(5.75), height=Inches(5.09))
        except Exception:
            missing_fields.append('SENTIMENT_PIE')
    else:
        missing_fields.append('SENTIMENT_PIE')

    # Wordcloud
    wc_added = False
    if wordcloud_path and os.path.exists(wordcloud_path):
        wc_added = place_image_at_placeholder('WORDCLOUD', wordcloud_path, default_size=(Inches(4.2), Inches(2.66)))
    if not wc_added:
        missing_fields.append('WORDCLOUD')

    # Top news (texto grande)
    topnews_slide, topnews_shape = find_shape_for_key('TOP_NEWS')
    if topnews_shape:
        try:
            set_text_style(topnews_shape, "\n".join(top_sentences), 'Effra Light', Pt(12), False)
        except Exception:
            missing_fields.append('TOP_NEWS')
    else:
        missing_fields.append('TOP_NEWS')

    # Análisis Groq
    analisis_texto = "No disponible"
    try:
        parrafos = "\n".join(df_cleaned['Hit Sentence'].dropna().astype(str).tolist()[:80])
        prompt = construir_prompt(client_name, parrafos)
        respuesta = llamar_groq(prompt)
        if respuesta:
            resultado_json = extraer_json(respuesta)
            if isinstance(resultado_json, dict):
                analisis_texto = formatear_analisis_social_listening(resultado_json)
    except Exception:
        analisis_texto = "No disponible"

    analisis_slide, analisis_shape = find_shape_for_key('CONVERSATION_ANALISIS')
    if analisis_shape:
        try:
            set_text_style(analisis_shape, analisis_texto, 'Effra Light', Pt(11), False)
        except Exception:
            missing_fields.append('CONVERSATION_ANALISIS')
    else:
        missing_fields.append('CONVERSATION_ANALISIS')

    # KPI NUMB_PRENSA / NUMB_REDES and tables: localizar placeholder y ubicar tabla
    prensa_shape_key = 'NUMB_PRENSA'
    prensa_slide, prensa_shape = find_shape_for_key(prensa_shape_key)
    if prensa_shape:
        try:
            set_text_style(prensa_shape, str(platform_counts.get('Prensa Digital', 0)), font_size=Pt(28))
        except Exception:
            missing_fields.append(prensa_shape_key)
    else:
        missing_fields.append(prensa_shape_key)

    try:
        table_added = False
        slide_for_table, shape_for_table = find_shape_for_key('TOP_INFLUENCERS_PRENSA_TABLE')
        if slide_for_table and shape_for_table:
            left, top, width, height = shape_for_table.left, shape_for_table.top, shape_for_table.width, shape_for_table.height
            try:
                ppt_engine.add_dataframe_as_table(slide_for_table, top_influencers_prensa, left, top, width, height)
                table_added = True
            except Exception:
                table_added = False
        if not table_added:
            missing_fields.append('TOP_INFLUENCERS_PRENSA_TABLE')
    except Exception:
        missing_fields.append('TOP_INFLUENCERS_PRENSA_TABLE')

    redes_shape_key = 'NUMB_REDES'
    redes_slide, redes_shape = find_shape_for_key(redes_shape_key)
    if redes_shape:
        try:
            set_text_style(redes_shape, str(platform_counts.get('Redes Sociales', 0)), font_size=Pt(28))
        except Exception:
            missing_fields.append(redes_shape_key)
    else:
        missing_fields.append(redes_shape_key)

    # Two tables for redes (posts and reach)
    try:
        table1_added = False
        slide_t1, shape_t1 = find_shape_for_key('TOP_INFLUENCERS_REDES_POSTS_TABLE')
        if slide_t1 and shape_t1:
            try:
                ppt_engine.add_dataframe_as_table(slide_t1, top_influencers_redes_posts, shape_t1.left, shape_t1.top, shape_t1.width, shape_t1.height)
                table1_added = True
            except Exception:
                table1_added = False
        if not table1_added:
            missing_fields.append('TOP_INFLUENCERS_REDES_POSTS_TABLE')
    except Exception:
        missing_fields.append('TOP_INFLUENCERS_REDES_POSTS_TABLE')

    try:
        table2_added = False
        slide_t2, shape_t2 = find_shape_for_key('TOP_INFLUENCERS_REDES_REACH_TABLE')
        if slide_t2 and shape_t2:
            try:
                ppt_engine.add_dataframe_as_table(slide_t2, top_influencers_redes_reach, shape_t2.left, shape_t2.top, shape_t2.width, shape_t2.height)
                table2_added = True
            except Exception:
                table2_added = False
        if not table2_added:
            missing_fields.append('TOP_INFLUENCERS_REDES_REACH_TABLE')
    except Exception:
        missing_fields.append('TOP_INFLUENCERS_REDES_REACH_TABLE')

    # Guardado de archivos
    safe_title = secure_filename(report_title) if report_title else f"Reporte_{unique_id}"
    pptx_filename = f"{safe_title}.pptx"
    pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], pptx_filename)
    prs.save(pptx_path)

    zip_filename = f"{safe_title}.zip"
    zip_path = os.path.join(app.config['UPLOAD_FOLDER'], zip_filename)
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        zipf.write(pptx_path, arcname=pptx_filename)
        zipf.write(processed_csv_path, arcname=os.path.basename(processed_csv_path))

    # Persistencia del reporte con plantilla usada
    new_report = Report(
        filename=zip_filename,
        user_id=current_user.id,
        title=report_title,
        description=description,
        template_name=template_filename
    )
    db.session.add(new_report)
    db.session.commit()

    return zip_path, missing_fields, template_filename


@app.route('/download/<path:filename>')
@login_required
def download_file(filename):
    # Path sanitization (S7)
    safe_filename = secure_filename(filename)
    full_path = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)
    
    # Verify path is within UPLOAD_FOLDER
    if not os.path.abspath(full_path).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
        app.logger.warning(f"Path traversal attempt: {filename}")
        abort(403)
    
    report_row = Report.query.filter_by(filename=safe_filename).first()
    if report_row is None:
        abort(404)

    if not current_user.is_admin and report_row.user_id != current_user.id:
        abort(403)

    if not os.path.exists(full_path):
        abort(404)

    @after_this_request
    def cleanup(response):
        clean_scratch_folder()
        return response

    log_activity('download_report', f'Descarga: {safe_filename}')
    return send_file(full_path, as_attachment=True)


@app.route('/mis-reportes')
@login_required
def mis_reportes():
    user_reports = Report.query.filter_by(user_id=current_user.id).order_by(Report.created_at.desc()).all()
    return render_template('mis_reportes.html', reports=user_reports)


@app.context_processor
def inject_current_year():
    return {'current_year': datetime.now().year}


@app.route('/error/archivo-invalido')
def error_archivo_invalido():
    return render_template('error.html',
                           title="Archivo inválido",
                           message="El archivo que subiste no es válido o tiene un formato incorrecto.")


@app.route('/clasificacion', methods=['GET', 'POST'])
@tool_required('classification')
def clasificacion():
    if request.method == 'POST':
        # 1. Obtener archivo, reglas y valor por defecto
        file = request.files.get('csv_file')
        rules_str = request.form.get('rules')
        default_val = request.form.get('default_val', 'Sin Clasificar')
        use_keywords = request.form.get('use_keywords') == 'true'
        
        if not file or not rules_str:
            flash("Archivo o reglas no proporcionados.", "error")
            return redirect(url_for('clasificacion'))
        
        try:
            rules = json.loads(rules_str)
        except json.JSONDecodeError:
            flash("Error al procesar las reglas de clasificación.", "error")
            return redirect(url_for('clasificacion'))
            
        # 2. Guardar archivo temporal
        unique_id = str(uuid.uuid4())
        filename = f"{unique_id}_{file.filename}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # 3. Clasificar
        try:
            print(f"DEBUG: Iniciando clasificación con default_val='{default_val}', use_keywords={use_keywords}")
            df_classified = classify_mentions(file_path, rules, default_val=default_val, use_keywords=use_keywords)
            
            # 4. Calcular Estadísticas (Distribución e Insights) - Optimized (P2)
            stats = {}
            if df_classified is not None and not df_classified.empty:
                # Vectorized approach: much faster than iterrows()
                grouped = df_classified.groupby(['Categoria', 'Tematica']).size().reset_index(name='count')
                
                for _, row in grouped.iterrows():
                    cat = str(row['Categoria'])
                    tem = str(row['Tematica'])
                    count = row['count']
                    
                    if cat not in stats:
                        stats[cat] = {"total": 0, "tematicas": {}}
                    
                    stats[cat]["total"] += count
                    stats[cat]["tematicas"][tem] = count
            
            # Insights adicionales para el visual
            top_category = "N/A"
            max_val = -1
            for cat, data in stats.items():
                if cat != default_val and data['total'] > max_val:
                    max_val = data['total']
                    top_category = cat

            # 5. Guardar resultado (CSV para descarga)
            output_filename = f"Clasificado_{file.filename}"
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"classified_{unique_id}.csv")
            df_classified.to_csv(output_path, sep='\t', encoding='utf-16', index=False)
            _register_temp_artifact('classified', unique_id, f"classified_{unique_id}.csv")

            log_activity('classify_data', f'Clasificación: {file.filename} ({len(df_classified)} filas, {len(stats)} categorías)')
            
            return jsonify({
                "success": True,
                "download_url": url_for('download_classified', file_id=unique_id, original_name=output_filename),
                "stats": stats,
                "total_rows": len(df_classified),
                "insights": {
                    "top_category": top_category,
                    "top_count": max_val if max_val != -1 else 0
                }
            })
            
        except Exception as e:
            app.logger.error(f"Error en clasificación: {e}")
            return jsonify({"success": False, "error": "Error procesando la clasificación. Por favor intenta nuevamente."}), 500
            
    return render_template('clasificacion.html')


# ─────────────────────────────────────────────────────────────


# ─────────────────────────────────────────────────────────────
# File format detection
# ─────────────────────────────────────────────────────────────

@app.route('/clasificacion/detect', methods=['POST'])
@tool_required('classification')
def clasificacion_detect():
    """Auto-detect format + return column list and preview rows."""
    file = request.files.get('csv_file')
    if not file:
        return jsonify({'success': False, 'error': 'No se recibio ningun archivo.'}), 400
    try:
        raw = file.read()
        result = detect_format(raw, file.filename)
        if result.get('error'):
            return jsonify({'success': False, 'error': result['error']}), 400
        return jsonify({'success': True,
                        'columns': result['columns'],
                        'preview': result['preview'],
                        'encoding': result.get('encoding'),
                        'sep': result.get('sep'),
                        'file_type': result.get('file_type')})
    except Exception as e:
        app.logger.error(f"Error en detect: {e}")
        return jsonify({'success': False, 'error': 'No se pudo analizar el archivo.'}), 500


# ─────────────────────────────────────────────────────────────
# Full-file upload for chunked classification
# ─────────────────────────────────────────────────────────────

@app.route('/clasificacion/upload', methods=['POST'])
@tool_required('classification')
def clasificacion_upload():
    """
    Receive the full file, read it properly (respecting encoding/sep overrides),
    convert to UTF-8 TSV, store in a session temp file, and return chunk metadata.
    This avoids using the browser's file.text() API which always decodes as UTF-8.
    """
    from services.file_loader import detect_format as _detect_fmt, read_full_as_tsv as _read_tsv

    file = request.files.get('csv_file')
    if not file:
        return jsonify({'success': False, 'error': 'No se recibio ningun archivo.'}), 400

    try:
        raw = file.read()

        # Auto-detect format first
        fmt = _detect_fmt(raw, file.filename)
        if fmt.get('error'):
            return jsonify({'success': False, 'error': fmt['error']}), 400

        # Apply manual overrides if provided
        manual_encoding = (request.form.get('encoding') or '').strip() or None
        manual_sep      = (request.form.get('sep') or '').strip() or None

        if manual_encoding:
            fmt['encoding'] = manual_encoding
        if manual_sep:
            fmt['sep'] = manual_sep

        # Read full file as UTF-8 TSV using the correct encoding/format
        header, body = _read_tsv(raw, fmt)
        if not header:
            return jsonify({'success': False, 'error': 'No se pudo leer el archivo con el formato indicado.'}), 400

        # Count data rows
        data_lines = [l for l in body.split('\n') if l.strip()]
        total_rows = len(data_lines)

        # Store TSV in a session temp file for the chunk endpoint to use
        session_id = uuid.uuid4().hex
        session_file = os.path.join(app.config['UPLOAD_FOLDER'], f"upload_{session_id}.tsv")
        with open(session_file, 'w', encoding='utf-8') as f:
            f.write(header)
            f.write(body)

        CHUNK_SIZE = 2000
        total_chunks = max(1, -(-total_rows // CHUNK_SIZE))  # ceiling division

        return jsonify({
            'success': True,
            'session_id': session_id,
            'header': header,
            'total_rows': total_rows,
            'total_chunks': total_chunks,
            'chunk_size': CHUNK_SIZE,
        })

    except Exception as e:
        app.logger.error(f"Error en clasificacion/upload: {e}")
        return jsonify({'success': False, 'error': 'Error leyendo el archivo. Verifica el formato y la codificacion.'}), 500


# ─────────────────────────────────────────────────────────────
# Serve TSV body for chunked classification
# ─────────────────────────────────────────────────────────────

@app.route('/clasificacion/upload_body/<session_id>', methods=['GET'])
@tool_required('classification')
def clasificacion_upload_body(session_id):
    """
    Return the body (data rows, no header) of the stored UTF-8 TSV session file
    so the frontend can split it into chunks without touching the raw binary file.
    """
    safe_sid = secure_filename(session_id)
    if not safe_sid or safe_sid != session_id:
        abort(400)
    session_file = os.path.join(app.config['UPLOAD_FOLDER'], f"upload_{safe_sid}.tsv")
    if not os.path.exists(session_file):
        abort(404)
    try:
        with open(session_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        # First line is the header; return only the body rows
        body = ''.join(lines[1:])
        from flask import Response
        return Response(body, mimetype='text/plain; charset=utf-8')
    except Exception as e:
        app.logger.error(f"Error sirviendo upload_body: {e}")
        abort(500)


# ─────────────────────────────────────────────────────────────
# Classification Presets CRUD
# ─────────────────────────────────────────────────────────────

@app.route('/clasificacion/presets', methods=['GET'])
@tool_required('classification')
def presets_list():
    presets = (ClassificationPreset.query
               .filter_by(user_id=current_user.id)
               .order_by(ClassificationPreset.created_at.desc())
               .all())
    return jsonify([{'id': p.id, 'name': p.name,
                     'created_at': p.created_at.strftime('%d/%m/%Y')} for p in presets])


@app.route('/clasificacion/presets', methods=['POST'])
@tool_required('classification')
def presets_create():
    data = request.get_json(force=True)
    name  = (data.get('name') or '').strip()[:100]
    rules = data.get('rules', [])
    if not name:
        return jsonify({'success': False, 'error': 'El nombre del preset es obligatorio.'}), 400
    try:
        preset = ClassificationPreset(
            user_id=current_user.id,
            name=name,
            rules_json=json.dumps(rules, ensure_ascii=False)
        )
        db.session.add(preset)
        db.session.commit()
        log_activity('preset_create', f'Preset guardado: {name}')
        return jsonify({'success': True, 'id': preset.id, 'name': preset.name})
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"Error guardando preset: {e}")
        return jsonify({'success': False, 'error': 'No se pudo guardar el preset.'}), 500


@app.route('/clasificacion/presets/<int:preset_id>', methods=['GET'])
@tool_required('classification')
def presets_load(preset_id):
    preset = ClassificationPreset.query.filter_by(id=preset_id, user_id=current_user.id).first()
    if not preset:
        return jsonify({'success': False, 'error': 'Preset no encontrado.'}), 404
    return jsonify({'success': True, 'rules': preset.get_rules(), 'name': preset.name})


@app.route('/clasificacion/presets/<int:preset_id>', methods=['DELETE'])
@tool_required('classification')
def presets_delete(preset_id):
    preset = ClassificationPreset.query.filter_by(id=preset_id, user_id=current_user.id).first()
    if not preset:
        return jsonify({'success': False, 'error': 'Preset no encontrado.'}), 404
    try:
        name = preset.name
        db.session.delete(preset)
        db.session.commit()
        log_activity('preset_delete', f'Preset eliminado: {name}')
        return jsonify({'success': True})
    except Exception as e:
        db.session.rollback()
        return jsonify({'success': False, 'error': 'No se pudo eliminar el preset.'}), 500


@app.route('/clasificacion/presets/<int:preset_id>', methods=['PUT'])
@tool_required('classification')
def presets_update(preset_id):
    preset = ClassificationPreset.query.filter_by(id=preset_id, user_id=current_user.id).first()
    if not preset:
        return jsonify({'success': False, 'error': 'Preset no encontrado.'}), 404
    try:
        data = request.get_json(force=True)
        new_name = (data.get('name') or '').strip()[:100]
        new_rules = data.get('rules')
        if new_name:
            preset.name = new_name
        if new_rules is not None:
            preset.rules_json = json.dumps(new_rules, ensure_ascii=False)
        db.session.commit()
        log_activity('preset_update', f'Preset actualizado: {preset.name}')
        return jsonify({'success': True, 'id': preset.id, 'name': preset.name})
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"Error actualizando preset: {e}")
        return jsonify({'success': False, 'error': 'No se pudo actualizar el preset.'}), 500

# Chunked classification endpoints
# ─────────────────────────────────────────────────────────────

@app.route('/clasificacion/chunk', methods=['POST'])
@tool_required('classification')
def clasificacion_chunk():
    """Receive one batch of CSV rows, classify it, append to a temp session file."""
    try:
        data = request.get_json(force=True)
        if not data:
            return jsonify({"success": False, "error": "No se recibieron datos."}), 400

        session_id   = data.get('session_id', '')
        header_text  = data.get('header', '')
        rows_text    = data.get('rows', '')
        rules        = data.get('rules', [])
        default_val  = data.get('default_val', 'Sin Clasificar')
        use_keywords = bool(data.get('use_keywords', False))
        chunk_index  = int(data.get('chunk_index', 0))
        text_col     = data.get('text_col', 'Hit Sentence') or 'Hit Sentence'
        keywords_col = data.get('keywords_col', '') or ''

        safe_sid = secure_filename(session_id)
        if not safe_sid or safe_sid != session_id:
            return jsonify({"success": False, "error": "session_id invalido."}), 400

        if not header_text or not rows_text:
            return jsonify({"success": False, "error": "Datos de chunk vacios."}), 400

        from services.classifier import classify_chunk as _classify_chunk
        df_chunk = _classify_chunk(rows_text, header_text, rules,
                                   default_val=default_val, use_keywords=use_keywords,
                                   text_col=text_col, keywords_col=keywords_col)

        if df_chunk is None or df_chunk.empty:
            return jsonify({"success": True, "partial_stats": {}, "rows_in_chunk": 0})

        session_file = os.path.join(app.config['UPLOAD_FOLDER'], f"session_{safe_sid}.csv")
        write_header = (chunk_index == 0) or (not os.path.exists(session_file))
        df_chunk.to_csv(session_file, sep='\t', encoding='utf-16',
                        index=False, mode='w' if write_header else 'a',
                        header=write_header)

        partial_stats = {}
        grouped = df_chunk.groupby(['Categoria', 'Tematica']).size().reset_index(name='count')
        for _, row in grouped.iterrows():
            cat   = str(row['Categoria'])
            tem   = str(row['Tematica'])
            count = int(row['count'])
            if cat not in partial_stats:
                partial_stats[cat] = {"total": 0, "tematicas": {}}
            partial_stats[cat]["total"] += count
            partial_stats[cat]["tematicas"][tem] = partial_stats[cat]["tematicas"].get(tem, 0) + count

        return jsonify({
            "success": True,
            "partial_stats": partial_stats,
            "rows_in_chunk": len(df_chunk)
        })

    except Exception as e:
        app.logger.error(f"Error en chunk de clasificacion: {e}")
        return jsonify({"success": False, "error": "Error procesando el chunk."}), 500


@app.route('/clasificacion/finalize', methods=['POST'])
@tool_required('classification')
def clasificacion_finalize():
    """Read the assembled session file, compute final stats, return download URL."""
    try:
        data = request.get_json(force=True)
        if not data:
            return jsonify({"success": False, "error": "No se recibieron datos."}), 400

        session_id    = data.get('session_id', '')
        original_name = data.get('original_name', 'archivo.csv')
        default_val   = data.get('default_val', 'Sin Clasificar')

        safe_sid = secure_filename(session_id)
        if not safe_sid or safe_sid != session_id:
            return jsonify({"success": False, "error": "session_id invalido."}), 400

        session_file = os.path.join(app.config['UPLOAD_FOLDER'], f"session_{safe_sid}.csv")
        if not os.path.exists(session_file):
            return jsonify({"success": False, "error": "Sesion no encontrada. Reinicia el proceso."}), 404

        df_full = pd.read_csv(session_file, sep='\t', encoding='utf-16', on_bad_lines='skip')

        stats = {}
        if not df_full.empty and 'Categoria' in df_full.columns and 'Tematica' in df_full.columns:
            grouped = df_full.groupby(['Categoria', 'Tematica']).size().reset_index(name='count')
            for _, row in grouped.iterrows():
                cat   = str(row['Categoria'])
                tem   = str(row['Tematica'])
                count = int(row['count'])
                if cat not in stats:
                    stats[cat] = {"total": 0, "tematicas": {}}
                stats[cat]["total"] += count
                stats[cat]["tematicas"][tem] = stats[cat]["tematicas"].get(tem, 0) + count

        top_category = "N/A"
        max_val = -1
        for cat, d in stats.items():
            if cat != default_val and d['total'] > max_val:
                max_val = d['total']
                top_category = cat

        safe_orig = secure_filename(original_name)
        output_filename = f"Clasificado_{safe_orig}"
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"classified_{safe_sid}.csv")
        os.replace(session_file, output_path)
        _register_temp_artifact('classified', safe_sid, f"classified_{safe_sid}.csv")

        log_activity('classify_data',
                     f'Clasificacion (chunked): {safe_orig} ({len(df_full)} filas, {len(stats)} categorias)')

        return jsonify({
            "success": True,
            "download_url": url_for('download_classified', file_id=safe_sid, original_name=output_filename),
            "stats": stats,
            "total_rows": len(df_full),
            "insights": {
                "top_category": top_category,
                "top_count": max_val if max_val != -1 else 0
            }
        })

    except Exception as e:
        app.logger.error(f"Error en finalizacion de clasificacion: {e}")
        return jsonify({"success": False, "error": "Error finalizando la clasificacion."}), 500

@app.route('/download_classified/<file_id>/<original_name>')
@login_required
def download_classified(file_id, original_name):
    # Path sanitization (S7)
    safe_id = secure_filename(file_id)
    safe_name = secure_filename(original_name)
    artifact = _get_owned_artifact_or_403('classified', safe_id)
    file_path = _scratch_path(artifact.storage_name)
    
    # Verify path is within UPLOAD_FOLDER
    if not os.path.abspath(file_path).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
        app.logger.warning(f"Path traversal attempt: {file_id}")
        abort(403)
    
    @after_this_request
    def cleanup(response):
        try:
            if os.path.exists(file_path):
                # Opcional: borrar archivo después de descarga
                # os.remove(file_path) 
                pass
        except Exception as e:
            app.logger.error(f"Error limpiando archivo clasificado: {e}")
        return response

    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True, download_name=safe_name)
    else:
        abort(404)


@app.route('/union')
@tool_required('file_merge')
def union_archivos():
    return render_template('union.html')


@app.route('/union/detect', methods=['POST'])
@tool_required('file_merge')
def union_detect():
    """Auto-detect format of an uploaded file and return columns + preview."""
    file = request.files.get('file')
    if not file:
        return jsonify({'success': False, 'error': 'No se recibio ningun archivo.'}), 400
    try:
        raw = file.read()
        result = detect_format(raw, file.filename)
        if result.get('error'):
            return jsonify({'success': False, 'error': result['error']}), 400
        return jsonify({
            'success': True,
            'columns': result['columns'],
            'preview': result['preview'],
            'encoding': result.get('encoding'),
            'sep': result.get('sep'),
            'file_type': result.get('file_type'),
        })
    except Exception as e:
        app.logger.error(f"Error en union/detect: {e}")
        return jsonify({'success': False, 'error': 'No se pudo analizar el archivo.'}), 500


@app.route('/union/merge', methods=['POST'])
@tool_required('file_merge')
def union_merge():
    """Merge uploaded files using default or advanced mode."""
    from services.file_merger import read_file, merge_default, merge_advanced, save_merged

    mode = request.form.get('mode', 'default')

    try:
        if mode == 'advanced':
            # Advanced: exactly 2 files + column mapping
            file_a = request.files.get('file_a')
            file_b = request.files.get('file_b')
            if not file_a or not file_b:
                return jsonify({'success': False, 'error': 'Se necesitan ambos archivos para el modo avanzado.'}), 400

            mapping_str = request.form.get('mapping', '{}')
            try:
                mapping = json.loads(mapping_str)
            except json.JSONDecodeError:
                return jsonify({'success': False, 'error': 'Mapeo de columnas invalido.'}), 400

            enc_a = request.form.get('encoding_a') or None
            sep_a = request.form.get('sep_a') or None
            enc_b = request.form.get('encoding_b') or None
            sep_b = request.form.get('sep_b') or None

            raw_a = file_a.read()
            raw_b = file_b.read()
            df_a = read_file(raw_a, file_a.filename, encoding=enc_a, sep=sep_a)
            df_b = read_file(raw_b, file_b.filename, encoding=enc_b, sep=sep_b)

            merged = merge_advanced(df_a, df_b, mapping)

            # Apply extra columns (if any)
            extra_cols_str = request.form.get('extra_columns', '[]')
            try:
                extra_cols = json.loads(extra_cols_str)
            except json.JSONDecodeError:
                extra_cols = []
            for ec in extra_cols:
                col_name = ec.get('name', '').strip()
                if col_name:
                    merged[col_name] = ec.get('value', '')

            files_merged = 2
            detail = f'Union avanzada: {file_a.filename} + {file_b.filename} ({len(merged)} filas)'

        else:
            # Default: 2+ files
            files = request.files.getlist('files')
            if len(files) < 2:
                return jsonify({'success': False, 'error': 'Se necesitan al menos 2 archivos.'}), 400

            # Per-file encoding/separator overrides come as JSON arrays
            encodings_str = request.form.get('encodings', '[]')
            seps_str = request.form.get('separators', '[]')
            try:
                encodings = json.loads(encodings_str)
                seps = json.loads(seps_str)
            except json.JSONDecodeError:
                encodings, seps = [], []

            dataframes = []
            filenames = []
            for i, f in enumerate(files):
                raw = f.read()
                enc = encodings[i] if i < len(encodings) and encodings[i] else None
                sep = seps[i] if i < len(seps) and seps[i] else None
                df = read_file(raw, f.filename, encoding=enc, sep=sep)
                dataframes.append(df)
                filenames.append(f.filename)

            merged = merge_default(dataframes)
            files_merged = len(files)
            detail = f'Union predeterminada: {", ".join(filenames)} ({len(merged)} filas)'

        # Save result
        unique_id = uuid.uuid4().hex[:10]
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"merged_{unique_id}.csv")
        save_merged(merged, output_path)
        _register_temp_artifact('union', unique_id, f"merged_{unique_id}.csv")

        log_activity('file_merge', detail)

        return jsonify({
            'success': True,
            'download_url': url_for('union_download', file_id=unique_id),
            'total_rows': len(merged),
            'total_columns': len(merged.columns),
            'files_merged': files_merged,
        })

    except ValueError as e:
        return jsonify({'success': False, 'error': str(e)}), 400
    except Exception as e:
        app.logger.error(f"Error en union/merge: {e}")
        return jsonify({'success': False, 'error': 'Error procesando la union de archivos.'}), 500


@app.route('/union/download/<file_id>')
@login_required
def union_download(file_id):
    """Download a merged file."""
    safe_id = secure_filename(file_id)
    artifact = _get_owned_artifact_or_403('union', safe_id)
    file_path = _scratch_path(artifact.storage_name)

    if not os.path.abspath(file_path).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
        abort(403)

    if not os.path.exists(file_path):
        abort(404)

    return send_file(file_path, as_attachment=True, download_name=f"Union_{safe_id}.csv")


@app.route('/analisis-csv', methods=['GET', 'POST'])
@tool_required('csv_analysis')
def analisis_csv():
    if request.method == 'POST':
        # Get file and parameters
        file = request.files.get('csv_file')
        encoding = request.form.get('encoding', 'utf-8')
        separator = request.form.get('separator', ',')
        
        if not file:
            return jsonify({'success': False, 'error': 'No se proporcionó ningún archivo'}), 400
        
        # Save file temporarily
        unique_id = str(uuid.uuid4())
        filename = f"{unique_id}_{secure_filename(file.filename)}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        try:
            # Run analysis
            result = analyze_csv(file_path, encoding, separator)

            if result['success']:
                log_activity('csv_analysis', f'Análisis CSV: {file.filename}')
                # Generate summary CSV for download
                summary_filename = f"summary_{unique_id}.csv"
                summary_path = os.path.join(app.config['UPLOAD_FOLDER'], summary_filename)
                generate_summary_csv(result, summary_path)
                _register_temp_artifact('csv_summary', unique_id, summary_filename)
                
                # Add download URL to result
                result['download_url'] = url_for('download_csv_summary', file_id=unique_id)
                result['original_filename'] = file.filename
                
                return jsonify(result)
            else:
                return jsonify(result), 400
                
        except Exception as e:
            app.logger.error(f"Error en análisis CSV: {e}")
            return jsonify({
                'success': False,
                'error': f'Error procesando el archivo: {str(e)}'
            }), 500
        finally:
            # Clean up uploaded file
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except Exception as e:
                    app.logger.error(f"Error eliminando archivo temporal: {e}")
    
    return render_template('analisis_csv.html')


@app.route('/analisis-csv/download/<file_id>')
@login_required
def download_csv_summary(file_id):
    # Path sanitization
    safe_id = secure_filename(file_id)
    artifact = _get_owned_artifact_or_403('csv_summary', safe_id)
    file_path = _scratch_path(artifact.storage_name)
    
    # Verify path is within UPLOAD_FOLDER
    if not os.path.abspath(file_path).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
        app.logger.warning(f"Path traversal attempt: {file_id}")
        abort(403)
    
    @after_this_request
    def cleanup(response):
        try:
            if os.path.exists(file_path):
                # Delete file after download
                os.remove(file_path)
        except Exception as e:
            app.logger.error(f"Error limpiando archivo de resumen: {e}")
        return response
    
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True, download_name=f"analisis_resumen_{safe_id}.csv")
    else:
        abort(404)

@app.route('/upload_csv', methods=['POST'])
@login_required
@tool_required('reports')
def upload_csv():
    if 'csv_file' not in request.files:
        flash('No se seleccionó ningún archivo')
        return redirect(url_for('index'))
    
    file = request.files['csv_file']
    
    if file.filename == '':
        flash('Nombre de archivo vacío')
        return redirect(url_for('index'))

    if file:
        try:
            upload_folder = app.config['UPLOAD_FOLDER']
            os.makedirs(upload_folder, exist_ok=True)

            safe_name = secure_filename(file.filename)
            if not safe_name:
                flash('Nombre de archivo inválido')
                return redirect(url_for('index'))

            file_path = os.path.join(upload_folder, safe_name)
            if not os.path.abspath(file_path).startswith(_scratch_root_abs()):
                app.logger.warning(f"Path traversal attempt in upload_csv: {file.filename}")
                abort(403)

            file.save(file_path)
            
            # Capturamos el título del formulario HTML también
            titulo = request.form.get('report_title', 'Mi Reporte')

            report_context = report.create_report_context(
                file_path, 
                report_title=titulo
            )
            return render_template('editor.html', context=report_context)
            
        except Exception as e:
            app.logger.error(f"ERROR: {e}")
            flash(f"Error procesando el archivo: {str(e)}")
            return redirect(url_for('index'))

    return redirect(url_for('index'))


# ─────────────────────────────────────────────────────────────
# Centralized Error Handlers (HTML for browser, JSON for AJAX)
# ─────────────────────────────────────────────────────────────

@app.errorhandler(400)
def bad_request(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "Solicitud invalida."}), 400
    return render_template('error.html',
                           title="Error 400 - Solicitud inválida",
                           message="La solicitud no es válida. Verifica los datos e intenta de nuevo."), 400


@app.errorhandler(403)
def forbidden(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "No tienes permiso para realizar esta accion."}), 403
    return render_template('error.html',
                           title="Error 403 - Acceso denegado",
                           message="No tienes permiso para acceder a este recurso."), 403


@app.errorhandler(404)
def page_not_found(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "Recurso no encontrado."}), 404
    return render_template('error.html',
                           title="Error 404 - Página no encontrada",
                           message="La página que estás buscando no existe."), 404


@app.errorhandler(405)
def method_not_allowed(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "Metodo no permitido."}), 405
    return render_template('error.html',
                           title="Error 405 - Método no permitido",
                           message="El método HTTP utilizado no está permitido para este recurso."), 405


@app.errorhandler(413)
def request_entity_too_large(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "El archivo es demasiado grande. El limite maximo es 200 MB."}), 413
    return render_template('error.html',
                           title="Archivo demasiado grande",
                           message="El archivo que subiste supera el limite de 200 MB."), 413


@app.errorhandler(429)
def rate_limit_exceeded(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "Demasiadas solicitudes. Intenta mas tarde."}), 429
    return render_template('error.html',
                           title="Error 429 - Demasiadas solicitudes",
                           message="Has realizado demasiadas solicitudes. Por favor espera unos minutos e intenta de nuevo."), 429


@app.errorhandler(500)
def internal_error(e):
    if _is_ajax():
        return jsonify({"success": False, "error": "Error interno del servidor."}), 500
    return render_template('error.html',
                           title="Error 500 - Problema del servidor",
                           message="Ocurrió un error inesperado. Por favor, intenta más tarde."), 500

@app.route('/generate_pptx', methods=['POST'])
@login_required
@tool_required('reports')
def generate_pptx_route():
    try:
        # 1. Recibir el JSON con los datos editados
        data = request.json
        if not data:
            return "No se recibieron datos JSON", 400

        # 2. Definir rutas
        template_path = os.path.join('powerpoints', 'Reporte_plantilla.pptx') 
        
        # Verificar que la plantilla existe
        if not os.path.exists(template_path):
            return f"Error: No se encuentra la plantilla en {template_path}", 500

        filename = f"Reporte_{data['meta']['client_name']}.pptx"
        output_path = os.path.join('scratch', filename)

        # 3. Llamar al motor de generación
        ppt_engine.generate_pptx(data, template_path, output_path)

        # 4. Enviar el archivo al usuario
        return send_file(output_path, as_attachment=True, download_name=filename)

    except Exception as e:
        app.logger.error(f"ERROR GENERANDO PPT: {e}")
        return "Error generando el reporte. Por favor intenta nuevamente.", 500

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    debug_mode = os.environ.get('FLASK_DEBUG', 'false').lower() == 'true'
    app.run(debug=debug_mode, host='0.0.0.0', port=port)
