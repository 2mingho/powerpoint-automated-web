"""
Native PowerPoint chart generation module.

Replaces matplotlib-based PNG chart images with native, editable PowerPoint
chart objects using python-pptx's CategoryChartData and add_chart APIs.
"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_native_line_chart(slide, shape, labels, values, series_name='Mentions', width=None, height=None):
    """
    Replace a text placeholder with a native PowerPoint line chart.
    """
    left, top = shape.left, shape.top
    final_width = width if width is not None else shape.width
    final_height = height if height is not None else shape.height

    # Remove the text placeholder
    sp = shape._element
    sp.getparent().remove(sp)

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(series_name, values)

    # Insert native chart (using LINE instead of LINE_MARKERS as per "line_marker=no")
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, final_width, final_height, chart_data
    )
    chart = chart_frame.chart

    # --- Styling ---
    chart.has_title = False
    chart.has_legend = False

    # Series styling: orange line, width=5pt, smoothed
    series = chart.series[0]
    series.format.line.color.rgb = RGBColor(0xFF, 0xA5, 0x00)  # orange
    series.format.line.width = Pt(5)
    series.smooth = True  # smoothed_line=yes

    # Category axis (x-axis) styling
    cat_axis = chart.category_axis
    cat_axis.has_major_gridlines = False
    cat_axis.tick_labels.font.size = Pt(8)
    cat_axis.tick_labels.font.name = 'Arial'

    # Value axis (y-axis) styling
    val_axis = chart.value_axis
    val_axis.has_major_gridlines = False  # gridlines=no
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.name = 'Arial'

    return chart


def add_native_pie_chart(slide, shape, sentiment_data, width=None, height=None):
    """
    Replace a text placeholder with a native PowerPoint pie chart.
    """
    left, top = shape.left, shape.top
    final_width = width if width is not None else shape.width
    final_height = height if height is not None else shape.height

    # Remove the text placeholder
    sp = shape._element
    sp.getparent().remove(sp)

    # Build chart data
    labels = [d['label'] for d in sentiment_data]
    values = [d['value'] for d in sentiment_data]
    colors = [d['color'] for d in sentiment_data]

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('Sentiment', values)

    # Insert native chart
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, final_width, final_height, chart_data
    )
    chart = chart_frame.chart

    # --- Styling ---
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = 'Arial'

    # Data labels: inside_end, % with 2 decimals
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    data_labels.number_format = '0.00%'
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = False
    data_labels.font.size = Pt(11)
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Apply custom colors
    series = chart.series[0]
    for idx, hex_color in enumerate(colors):
        point = series.points[idx]
        hex_str = hex_color.lstrip('#')
        r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor(r, g, b)

    return chart
