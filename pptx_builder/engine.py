import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from . import native_charts

# Configuración de estilo
FONT_NAME = 'Arial' 

def set_text_style(shape, text, font_name=None, font_size=Pt(14), bold=False, color=RGBColor(0,0,0), alignment=PP_ALIGN.LEFT):
    """Helper para aplicar estilos a cuadros de texto"""
    if not shape.has_text_frame: return
    text_frame = shape.text_frame
    text_frame.clear() 
    p = text_frame.paragraphs[0]
    p.text = str(text)
    p.font.name = font_name if font_name else FONT_NAME
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

def add_dataframe_as_table(slide, data_or_shape, headers_or_data=None, left=None, top=None, width=None, height=None, headers=None):
    """
    Add a DataFrame or list-of-dicts as a native PowerPoint table.
    
    Supports two calling patterns:
    1. (slide, shape, data_list, headers) — replaces a placeholder shape
    2. (slide, data_list, left, top, width, height) — positions explicitly (legacy app.py pattern)
    """
    # Determine calling pattern
    if hasattr(data_or_shape, 'left'):
        # Pattern 1: shape-based — extract position from shape, remove it
        shape = data_or_shape
        data_list = headers_or_data if headers_or_data is not None else []
        final_left, final_top = shape.left, shape.top
        final_width = width if width else shape.width
        final_height = height if height else shape.height
        # Remove placeholder
        sp = shape._element
        sp.getparent().remove(sp)
    else:
        # Pattern 2: explicit positioning (legacy app.py pattern)
        # data_or_shape is actually the data (DataFrame or list)
        data_list = data_or_shape
        if headers is None and isinstance(headers_or_data, (int, float)):
            # Called as (slide, data, left, top, width, height)
            final_left = headers_or_data  # Actually 'left'
            final_top = left              # Actually 'top'  
            final_width = top             # Actually 'width'
            final_height = width          # Actually 'height'
            headers = None
        else:
            final_left = left
            final_top = top
            final_width = width
            final_height = height

    # Convert DataFrame to list of dicts if needed
    if hasattr(data_list, 'to_dict'):
        if headers is None:
            headers = list(data_list.columns)
        data_list = data_list.to_dict(orient='records')
    
    if headers is None:
        headers = list(data_list[0].keys()) if data_list else []

    # Standard slide dimensions (13.333" x 7.5" for widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Override with requested dimensions and center
    final_width = Inches(7.88)
    final_height = Inches(4.07)
    final_left = int((SLIDE_WIDTH - final_width) / 2)
    final_top = int((SLIDE_HEIGHT - final_height) / 2)

    rows = len(data_list) + 1
    cols = len(headers)
    if cols == 0:
        return

    table = slide.shapes.add_table(rows, cols, final_left, final_top, final_width, final_height).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 192, 0)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data
    for row_idx, row_data in enumerate(data_list):
        for col_idx, header in enumerate(headers):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(row_data.get(header, ''))
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def generate_pptx(json_data, template_path, output_path):
    """Función Maestra V3.1 (Corrección de argumentos)"""
    
    prs = Presentation(template_path)

    # --- 1. PREPARAR DATOS DE TEXTO ---
    top_news_text = "\n\n".join(json_data['tables']['top_sentences']) if json_data['tables']['top_sentences'] else "No hay noticias destacadas."
    analisis_text = json_data.get('ai_analysis', {}).get('summary', "Análisis no disponible.")

    text_replacements = {
        "REPORT_CLIENT": json_data['meta']['client_name'],
        "REPORT_DATE": json_data['meta']['date_generated'],
        "NUMB_MENTIONS": str(json_data['kpis']['total_mentions']),
        "NUMB_ACTORS": str(json_data['kpis']['unique_authors']),
        "EST_REACH": json_data['kpis']['estimated_reach_fmt'],
        "NUMB_PRENSA": str(json_data['kpis']['mentions_prensa']),
        "NUMB_REDES": str(json_data['kpis']['mentions_redes']),
        "TOP_NEWS": top_news_text,
        "CONVERSATION_ANALISIS": analisis_text
    }

    # --- 2. REEMPLAZO DE TEXTO ---
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for key, value in text_replacements.items():
                    if key in shape.text:
                        if key in ["TOP_NEWS", "CONVERSATION_ANALISIS"]:
                             set_text_style(shape, value, font_size=Pt(11), bold=False, alignment=PP_ALIGN.LEFT)
                        elif key.startswith("NUMB_") or key == "EST_REACH":
                             set_text_style(shape, value, font_size=Pt(28), bold=True, alignment=PP_ALIGN.CENTER)
                        elif key == "REPORT_DATE":
                             set_text_style(shape, value, font_size=Pt(24), bold=True, color=RGBColor(255, 255, 255), alignment=PP_ALIGN.CENTER)
                        else:
                             set_text_style(shape, value, font_size=Pt(24), bold=True, alignment=PP_ALIGN.CENTER)

    # --- 3. INSERCIÓN DE GRÁFICOS NATIVOS ---
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                txt = shape.text.strip()
                if txt == 'SENTIMENT_PIE':
                    native_charts.add_native_pie_chart(slide, shape, json_data['charts']['sentiment'], width=Inches(5.75), height=Inches(5.09))
                elif txt == 'CONVERSATION_CHART':
                    evo = json_data['charts']['evolution']
                    native_charts.add_native_line_chart(slide, shape, evo['labels'], evo['values'], width=Inches(9.07), height=Inches(5.15))

    # --- 4. TABLAS ---
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                txt = shape.text
                if 'TOP_INFLUENCERS_PRENSA_TABLE' in txt:
                    add_dataframe_as_table(slide, shape, json_data['tables']['top_prensa'], headers=['Influencer', 'Posts', 'Reach'])
                elif 'TOP_INFLUENCERS_REDES_POSTS_TABLE' in txt:
                    add_dataframe_as_table(slide, shape, json_data['tables']['top_redes'], headers=['Influencer', 'Posts', 'Reach', 'Source'])
                elif 'TOP_INFLUENCERS_REDES_REACH_TABLE' in txt:
                    add_dataframe_as_table(slide, shape, json_data['tables']['top_redes'], headers=['Influencer', 'Posts', 'Reach', 'Source'])

    prs.save(output_path)
    return output_path