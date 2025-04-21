import os
import textwrap
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ─────────────────────────────────────────────────────────────
social_network_sources = {
    'Twitter': 'Redes Sociales',
    'Youtube': 'Redes Sociales',
    'Instagram': 'Redes Sociales',
    'Facebook': 'Redes Sociales',
    'Pinterest': 'Redes Sociales',
    'Reddit': 'Redes Sociales',
    'TikTok': 'Redes Sociales',
    'Twitch': 'Redes Sociales',
}

def format_number(number):
    if number >= 1_000_000:
        return f"{number / 1_000_000:.1f}M"
    elif number >= 1_000:
        return f"{number / 1_000:.1f}k"
    return str(number)

def format_reach(number):
    if pd.isna(number):
        return number
    return f"{int(number):,}"

def set_text_style(shape, text, font_name='Effra', font_size=Pt(14), center=True):
    if shape.has_text_frame:
        p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
        p.clear()
        p.text = text
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

# ─────────────────────────────────────────────────────────────
def load_and_clean_data(file_path):
    df = pd.read_csv(file_path, encoding='utf-16', sep='\t')
    columns_to_delete = [
        'Opening Text', 'Subregion', 'Desktop Reach', 'Mobile Reach',
        'Twitter Social Echo', 'Facebook Social Echo', 'Reddit Social Echo',
        'National Viewership', 'AVE', 'State', 'City',
        'Social Echo Total', 'Editorial Echo', 'Views', 'Estimated Views',
        'Likes', 'Replies', 'Retweets', 'Comments', 'Shares', 'Reactions',
        'Threads', 'Is Verified'
    ]
    df_cleaned = df.drop(columns=columns_to_delete, errors='ignore')
    df_cleaned['Hit Sentence'] = df_cleaned['Headline'].where(~df_cleaned['Headline'].isna(), df_cleaned['Hit Sentence'])
    exclude_sources = list(social_network_sources.keys())
    mask = ~df_cleaned['Source'].isin(exclude_sources)
    df_cleaned.loc[mask, 'Influencer'] = df_cleaned.loc[mask, 'Source']
    df_cleaned['Plataforma'] = df_cleaned['Source'].apply(lambda x: social_network_sources.get(x, 'Prensa Digital'))
    return df_cleaned

def update_influencer(row):
    if row['Source'] == 'Facebook' and (row['Reach'] == 0 or pd.isna(row['Reach']) or row['Reach'] == ""):
        return "Comment on " + row['Influencer']
    return row['Influencer']

def update_sentiment(row):
    sentiment = row['Sentiment']
    return "Neutral" if sentiment == "Unknown" or pd.isna(sentiment) else sentiment

def calculate_summary_metrics(df):
    total_mentions = len(df)
    count_of_authors = df['Influencer'].nunique()
    estimated_reach = df.groupby('Influencer')['Reach'].max().sum()
    formatted_estimated_reach = format_number(estimated_reach)
    return total_mentions, count_of_authors, formatted_estimated_reach

# ─────────────────────────────────────────────────────────────
def create_mentions_evolution_chart(df, date_column='Alternate Date Format', time_column='Time', output_path='scratch/convEvolution.png'):
    df['date'] = pd.to_datetime(df[date_column], format='%d-%b-%y').dt.date
    df['hour'] = pd.to_datetime(df[time_column], format='%I:%M %p').dt.strftime('%I %p')
    df_grouped = df.groupby(['date', 'hour']).size().reset_index(name='count')
    df_grouped['datetime'] = pd.to_datetime(df_grouped['date'].astype(str) + ' ' + df_grouped['hour'])
    df_grouped = df_grouped.sort_values('datetime')

    fig, ax = plt.subplots(figsize=(13, 7))
    ax.plot(df_grouped['datetime'], df_grouped['count'], linewidth=3.5, color='orange')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.xaxis.set_major_locator(mdates.HourLocator(interval=2))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%d-%b %I %p'))
    plt.xticks(rotation=90, ha='right')
    plt.tight_layout()
    plt.savefig(output_path, transparent=True)
    plt.close(fig)

def create_mentions_evolution_chart_by_date(df, date_column='Alternate Date Format', output_path='scratch/convEvolution.png'):
    df['date'] = pd.to_datetime(df[date_column], format='%d-%b-%y').dt.date
    df_grouped = df.groupby('date').size().reset_index(name='count')
    df_grouped = df_grouped.sort_values('date')

    fig, ax = plt.subplots(figsize=(13, 7))
    ax.plot(df_grouped['date'], df_grouped['count'], linewidth=3.5, color='orange')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.xaxis.set_major_locator(mdates.DayLocator(interval=1))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%d-%b'))
    plt.xticks(rotation=90, ha='right')
    plt.tight_layout()
    plt.savefig(output_path, transparent=True)
    plt.close(fig)

def create_sentiment_pie_chart(df, output_path='scratch/sentiment_pie_chart.png'):
    sentiment_counts = df[df['Sentiment'] != 'Not Rated']['Sentiment'].value_counts()
    sentiment_colors = {'Negative': '#ff0000', 'Positive': '#00b050', 'Neutral': '#BFBFBF'}
    labels = sentiment_counts.index.to_list()
    sizes = sentiment_counts.to_list()

    fig, ax = plt.subplots(figsize=(10, 10))
    colors = [sentiment_colors.get(label, 'lightgray') for label in labels]
    ax.pie(sizes, autopct="%1.1f%%", colors=colors, textprops=dict(backgroundcolor='w', fontsize=14, fontweight='bold'), startangle=140)
    ax.axis('equal')
    plt.savefig(output_path, transparent=True)
    plt.close(fig)

# ─────────────────────────────────────────────────────────────
def distribucion_plataforma(df):
    platform_counts = df['Plataforma'].value_counts()
    max_reach_per_platform = df.groupby('Plataforma')['Reach'].max().apply(format_number)
    return platform_counts.to_dict(), max_reach_per_platform.to_dict()

def get_top_hit_sentences(df):
    top_influencers = df[df['Plataforma'] == "Prensa Digital"].sort_values(by='Reach', ascending=False).head(5)['Hit Sentence']
    return [textwrap.fill(sentence, width=100) for sentence in top_influencers]

def get_top_influencers(df, plataforma, sort_by='Posts', top_n=10, include_source=False):
    df_filtered = df[df['Plataforma'] == plataforma]

    if include_source:
        grouped = (
            df_filtered.groupby('Influencer')[['Reach', 'Source']]
            .agg(Posts=('Reach', 'count'), Max_Reach=('Reach', 'max'), Source=('Source', 'first'))
        )
        grouped.reset_index(inplace=True)
        grouped['Max_Reach'] = grouped['Max_Reach'].apply(format_reach)
        grouped = grouped.sort_values(by=sort_by, ascending=False).head(top_n)
        return grouped[['Influencer', 'Posts', 'Max_Reach', 'Source']]
    else:
        grouped = df_filtered.groupby('Influencer')['Reach'].agg(['count', 'max']).reset_index()
        grouped.columns = ['Influencer', 'Posts', 'Max Reach']
        grouped['Max Reach'] = grouped['Max Reach'].apply(format_reach)
        grouped = grouped.sort_values(by=sort_by, ascending=False).head(top_n)
        return grouped

# ─────────────────────────────────────────────────────────────
def add_dataframe_as_table(slide, dataframe, x, y, width, height):
    if dataframe.empty:
        print("DataFrame está vacío.")
        return

    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, x, y, width, height).table

    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = str(dataframe.columns[j])
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 192, 0)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.name = 'Effra'
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i in range(rows):
        for j in range(cols):
            value = dataframe.iat[i, j]
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Effra Light'
                paragraph.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    print("Tabla añadida correctamente.")

# ─────────────────────────────────────────────────────────────
def save_cleaned_csv(df, file_path, unique_id=None):
    base_filename = os.path.basename(file_path).split()[0].split('.')[0]
    if unique_id:
        output_filename = os.path.join("scratch", f"{base_filename}_{unique_id}_(resultado).csv")
    else:
        output_filename = os.path.join("scratch", f"{base_filename}_(resultado).csv")

    with open(output_filename, 'w', encoding='utf-16', newline='') as f:
        df.to_csv(f, index=False, sep='\t')

    print(f"CSV data saved to: {output_filename}")
    return output_filename