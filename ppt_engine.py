import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# Configuración de estilo
FONT_NAME = 'Arial' 

def set_text_style(shape, text, font_size=Pt(14), bold=False, color=RGBColor(0,0,0), alignment=PP_ALIGN.LEFT):
    """Helper para aplicar estilos a cuadros de texto"""
    if not shape.has_text_frame: return
    text_frame = shape.text_frame
    text_frame.clear() 
    p = text_frame.paragraphs[0]
    p.text = str(text)
    p.font.name = FONT_NAME
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

def create_temp_charts(data):
    """Regenera los gráficos con dimensiones optimizadas"""
    paths = {}
    os.makedirs('scratch', exist_ok=True)

    # 1. Gráfico de Sentimiento (Pastel)
    sent_data = data['charts']['sentiment']
    order = ['Positive', 'Neutral', 'Negative']
    sent_data.sort(key=lambda x: order.index(x['label']) if x['label'] in order else 99)

    labels = [d['label'] for d in sent_data]
    values = [d['value'] for d in sent_data]
    colors = [d['color'] for d in sent_data]

    # Ratio 1:1
    fig, ax = plt.subplots(figsize=(5, 5))
    wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors, startangle=90)
    
    plt.setp(autotexts, size=12, weight="bold", color="white")
    plt.setp(texts, size=10)
    
    sent_path = 'scratch/temp_sentiment.png'
    plt.savefig(sent_path, transparent=True, dpi=150, bbox_inches='tight')
    plt.close(fig)
    paths['sentiment'] = sent_path

    # 2. Gráfico de Evolución (Línea)
    evo_data = data['charts']['evolution']
    
    # Ratio 16:9 aprox (8.9 x 5)
    fig, ax = plt.subplots(figsize=(8.9, 5))
    ax.plot(evo_data['labels'], evo_data['values'], color='orange', linewidth=3, marker='o', markersize=6)
    
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_linewidth(0.5)
    ax.spines['bottom'].set_linewidth(0.5)
    ax.tick_params(axis='both', which='major', labelsize=10)
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    
    evo_path = 'scratch/temp_evolution.png'
    plt.savefig(evo_path, transparent=True, dpi=150)
    plt.close(fig)
    paths['evolution'] = evo_path

    return paths

def generate_pptx(json_data, template_path, output_path):
    """Función Maestra V3.1 (Corrección de argumentos)"""
    
    prs = Presentation(template_path)
    chart_images = create_temp_charts(json_data)

    # --- 1. PREPARAR DATOS DE TEXTO ---
    top_news_text = "\n\n".join(json_data['tables']['top_sentences']) if json_data['tables']['top_sentences'] else "No hay noticias destacadas."
    analisis_text = json_data.get('ai_analysis', {}).get('summary', "Análisis no disponible.")

    text_replacements = {
        "REPORT_CLIENT": json_data['meta']['client_name'],
        "REPORT_DATE": json_data['meta']['date_generated'],
        "NUMB_MENTIONS": str(json_data['kpis']['total_mentions']),
        "NUMB_ACTORS": str(json_data['kpis']['unique_authors']),
        "EST_REACH": json_data['kpis']['estimated_reach_fmt'],
        "NUMB_PRENSA": str(json_data['kpis']['mentions_prensa']),
        "NUMB_REDES": str(json_data['kpis']['mentions_redes']),
        "TOP_NEWS": top_news_text,
        "CONVERSATION_ANALISIS": analisis_text
    }

    # --- 2. REEMPLAZO DE TEXTO ---
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for key, value in text_replacements.items():
                    if key in shape.text:
                        if key in ["TOP_NEWS", "CONVERSATION_ANALISIS"]:
                             set_text_style(shape, value, font_size=Pt(11), bold=False, alignment=PP_ALIGN.LEFT)
                        elif key.startswith("NUMB_") or key == "EST_REACH":
                             set_text_style(shape, value, font_size=Pt(28), bold=True, alignment=PP_ALIGN.CENTER)
                        else:
                             set_text_style(shape, value, font_size=Pt(24), bold=True, alignment=PP_ALIGN.CENTER)

    # --- 3. INSERCIÓN DE IMÁGENES (CORREGIDO) ---
    # Ahora los argumentos coinciden: width y height
    def replace_image_placeholder(keyword, image_path, width=None, height=None):
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and keyword in shape.text:
                    left, top = shape.left, shape.top
                    
                    # Usamos dimensiones forzadas si existen, si no, las del placeholder
                    final_width = width if width else shape.width
                    final_height = height if height else shape.height
                    
                    shape.text = "" 
                    slide.shapes.add_picture(image_path, left, top, width=final_width, height=final_height)
                    return 

    # Aplicamos las medidas solicitadas
    replace_image_placeholder("SENTIMENT_PIE", chart_images['sentiment'], width=Inches(4.91), height=Inches(5))
    replace_image_placeholder("CONVERSATION_CHART", chart_images['evolution'], width=Inches(8.9), height=Inches(5))

    # --- 4. TABLAS (CORREGIDO) ---
    # Ahora los argumentos coinciden: width y height
    def create_native_table(slide, shape, data_list, headers, width=None, height=None):
        left, top = shape.left, shape.top
        
        # Dimensiones forzadas
        final_width = width if width else shape.width
        final_height = height if height else shape.height
        
        sp = shape._element
        sp.getparent().remove(sp)
        
        rows = len(data_list) + 1
        cols = len(headers)
        table = slide.shapes.add_table(rows, cols, left, top, final_width, final_height).table
        
        # Headers
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 192, 0)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Data
        for row_idx, row_data in enumerate(data_list):
            values_ordered = []
            if 'Source' in headers:
                 values_ordered = [row_data.get('Influencer'), row_data.get('Posts'), row_data.get('Reach'), row_data.get('Source')]
            else:
                 values_ordered = [row_data.get('Influencer'), row_data.get('Posts'), row_data.get('Reach')]

            for col_idx, value in enumerate(values_ordered):
                if col_idx < cols:
                    cell = table.cell(row_idx+1, col_idx)
                    cell.text = str(value)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Medidas solicitadas para tablas
    TABLE_WIDTH = Inches(5.43)
    TABLE_HEIGHT = Inches(4)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                txt = shape.text
                if 'TOP_INFLUENCERS_PRENSA_TABLE' in txt:
                    create_native_table(slide, shape, json_data['tables']['top_prensa'], ['Influencer', 'Posts', 'Reach'], width=TABLE_WIDTH, height=TABLE_HEIGHT)
                elif 'TOP_INFLUENCERS_REDES_POSTS_TABLE' in txt:
                    create_native_table(slide, shape, json_data['tables']['top_redes'], ['Influencer', 'Posts', 'Reach', 'Source'], width=TABLE_WIDTH, height=TABLE_HEIGHT)
                elif 'TOP_INFLUENCERS_REDES_REACH_TABLE' in txt:
                    create_native_table(slide, shape, json_data['tables']['top_redes'], ['Influencer', 'Posts', 'Reach', 'Source'], width=TABLE_WIDTH, height=TABLE_HEIGHT)

    prs.save(output_path)
    return output_path