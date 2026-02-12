import os
import sys
from pptx import Presentation

# Add root to path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

def run_test():
    """
    Test 5: PowerPoint Engine
    Returns: (success, details, data_points)
    """
    results = []
    template_path = "powerpoints/Reporte_plantilla.pptx"
    output_path = "scratch/test_output.pptx"
    
    try:
        # 1. Check template existence
        t_exists = os.path.exists(template_path)
        results.append({"name": "Plantilla Base", "status": t_exists, "diag": f"Falta {template_path}"})
        
        if not t_exists:
            return False, "No se puede probar el motor PPT sin la plantilla base.", {"Detalle": results}

        # 2. Try loading and saving a simple change
        prs = Presentation(template_path)
        # Add a test slide
        slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None
        if slide_layout:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            if title:
                title.text = "Test Diagnóstico"
        
        prs.save(output_path)
        
        s_ok = os.path.exists(output_path)
        results.append({"name": "Generación de PPTX", "status": s_ok, "diag": "Error al guardar archivo PPTX en scratch/"})
        
        if s_ok:
            os.remove(output_path)

        success = all(r['status'] for r in results)
        details = "El motor de PowerPoint funciona correctamente con la plantilla base." if success else "Error en la manipulación de presentaciones."
        
        data_points = {
            "Plantilla": template_path,
            "Slides Detectadas": len(prs.slides),
            "Detalle": results
        }

        return success, details, data_points

    except Exception as e:
        return False, f"Error en motor PPT: {str(e)}", {"Detalle": results}

if __name__ == "__main__":
    print(run_test())
